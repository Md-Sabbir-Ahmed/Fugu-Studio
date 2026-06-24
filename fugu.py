#!/usr/bin/env python3
"""
Fugu Studio — a desktop app for the Sakana Fugu / Fugu Ultra Responses API.

The Fugu API natively accepts TEXT and IMAGE input only (per the official
configs/files/fugu.json -> "input_modalities": ["text", "image"]).
So ANY file type can be attached/scanned — contents are handled locally:

  * image (png/jpg/webp/gif/bmp) -> sent as native input_image
  * pdf                          -> text extracted; scanned pages rendered to images
  * docx                         -> text extracted (paragraphs + tables)
  * pptx                         -> text extracted (all slides)
  * xlsx/xlsm                    -> text extracted (all sheets)
  * txt/md/csv/json/code/...     -> sent as text (any text-like file)
  * anything else (binary)       -> sent as a name/type/size placeholder

Run:  double-click this file, or `python fugu.py`
Deps: pip install openai pypdf python-docx python-pptx openpyxl pillow pymupdf tkinterdnd2
"""

import base64
import io
import json
import os
import re
import tempfile
import threading
import time
import traceback

# ---------------------------------------------------------------- credentials
# No key is baked in. The key comes from the env var or what you type in the
# app (persisted to api_key.txt). An empty key means "not configured".
_HERE = os.path.dirname(os.path.abspath(__file__))
_KEY_FILE = os.path.join(_HERE, "api_key.txt")


def load_key() -> str:
    """Priority: env var -> saved api_key.txt -> empty (not configured)."""
    key = os.environ.get("FUGU_API_KEY", "").strip()
    if key:
        return key
    try:
        with open(_KEY_FILE, "r", encoding="utf-8") as fh:
            return fh.read().strip()
    except FileNotFoundError:
        return ""


def set_api_key(key: str) -> None:
    """Update the in-memory key and persist it. Empty clears the saved file."""
    global API_KEY
    API_KEY = key.strip()
    try:
        if API_KEY:
            with open(_KEY_FILE, "w", encoding="utf-8") as fh:
                fh.write(API_KEY)
        elif os.path.exists(_KEY_FILE):
            os.remove(_KEY_FILE)
    except OSError:
        pass


API_KEY = load_key()
BASE_URL = os.environ.get("FUGU_BASE_URL", "https://api.sakana.ai").rstrip("/")
if not BASE_URL.endswith("/v1"):
    BASE_URL += "/v1"

MODELS = ["fugu", "fugu-ultra"]

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
TEXT_EXTS = {".txt", ".md", ".markdown", ".csv", ".json", ".log", ".py",
             ".js", ".ts", ".html", ".css", ".xml", ".yaml", ".yml", ".rtf",
             # extra source/code types so a project folder is fully recognised
             ".jsx", ".tsx", ".vue", ".svelte", ".mjs", ".cjs",
             ".java", ".kt", ".kts", ".c", ".h", ".cpp", ".cc", ".hpp",
             ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".m", ".mm",
             ".sh", ".bash", ".zsh", ".ps1", ".bat", ".cmd",
             ".sql", ".toml", ".ini", ".cfg", ".conf", ".env", ".properties",
             ".scss", ".sass", ".less", ".gradle", ".dockerfile",
             ".r", ".lua", ".pl", ".dart", ".ipynb", ".tex"}
MAX_IMAGE_DIM = 2000          # downscale larger images before upload
MAX_PDF_RENDER_PAGES = 15     # cap pages rendered for scanned PDFs
MAX_TEXT_CHARS = 200_000      # safety cap on extracted document text

# ----- folder ("workspace") scanning -------------------------------------
# Directories that are never worth sending (vcs, deps, build artefacts, caches).
IGNORE_DIRS = {".git", ".hg", ".svn", "node_modules", "__pycache__", ".venv",
               "venv", "env", ".env", "dist", "build", ".idea", ".vscode",
               ".mypy_cache", ".pytest_cache", ".ruff_cache", "site-packages",
               ".next", ".nuxt", "target", "bin", "obj", "out", ".gradle",
               ".cache", "coverage", ".parcel-cache", "vendor", ".tox", "chats"}
MAX_FOLDER_FILE_BYTES = 1_000_000     # skip any single file bigger than this
MAX_FOLDER_FILES = 400                # cap on number of files scanned
MAX_FOLDER_TOTAL_BYTES = 8_000_000    # cap on combined size scanned


def human_size(n: float) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if n < 1024 or unit == "GB":
            return f"{n:.0f} {unit}" if unit == "B" else f"{n:.1f} {unit}"
        n /= 1024
    return f"{n:.1f} GB"


def scan_folder(root_dir: str):
    """Walk root_dir, return (abs_paths, skipped_count). Smart-ignores junk."""
    root_dir = os.path.abspath(root_dir)
    found, total_bytes, skipped = [], 0, 0
    for dirpath, dirnames, filenames in os.walk(root_dir):
        # prune ignored / hidden dirs in place so os.walk skips them
        dirnames[:] = [d for d in dirnames
                       if d not in IGNORE_DIRS and not d.startswith(".")]
        for fn in sorted(filenames):
            # Accept EVERY file type — text is read, binaries become placeholders.
            path = os.path.join(dirpath, fn)
            try:
                size = os.path.getsize(path)
            except OSError:
                continue
            if size > MAX_FOLDER_FILE_BYTES:
                skipped += 1
                continue
            if len(found) >= MAX_FOLDER_FILES or total_bytes + size > MAX_FOLDER_TOTAL_BYTES:
                skipped += 1
                continue
            found.append(path)
            total_bytes += size
    found.sort(key=lambda p: os.path.relpath(p, root_dir).lower())
    return found, skipped


# ============================================================== file handling
def encode_image_bytes(raw: bytes) -> str:
    """Downscale if huge, return a data: URL."""
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(raw))
        if max(img.size) > MAX_IMAGE_DIM:
            img.thumbnail((MAX_IMAGE_DIM, MAX_IMAGE_DIM))
        buf = io.BytesIO()
        if img.mode in ("RGBA", "P", "LA"):
            img.save(buf, format="PNG"); mime = "image/png"
        else:
            img.convert("RGB").save(buf, format="JPEG", quality=85); mime = "image/jpeg"
        data = buf.getvalue()
    except Exception:
        data, mime = raw, "image/png"
    return f"data:{mime};base64," + base64.b64encode(data).decode()


def image_file_to_url(path: str) -> str:
    with open(path, "rb") as fh:
        return encode_image_bytes(fh.read())


def extract_pdf(path: str):
    """Return (text, [image_data_urls]). Renders pages to images if scanned."""
    text, images = "", []
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        parts = []
        for page in reader.pages:
            parts.append(page.extract_text() or "")
        text = "\n".join(parts).strip()
    except Exception as e:
        text = f"[pdf text extraction failed: {e}]"

    # If almost no text came out, it's likely scanned -> render pages as images.
    if len(text) < 40:
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(path)
            for i, page in enumerate(doc):
                if i >= MAX_PDF_RENDER_PAGES:
                    break
                pix = page.get_pixmap(dpi=150)
                images.append(encode_image_bytes(pix.tobytes("png")))
            doc.close()
            if images:
                text = ""  # rely on the rendered images
        except Exception as e:
            text += f"\n[pdf render failed: {e}]"
    return text, images


def extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def extract_xlsx(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"--- Sheet: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
    wb.close()
    return "\n".join(parts).strip()


def read_text_file(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="replace") as fh:
        return fh.read()


def looks_like_text(path: str, blocksize: int = 4096) -> bool:
    """Heuristic: True if the first chunk decodes cleanly as text (no NUL bytes)."""
    try:
        with open(path, "rb") as fh:
            chunk = fh.read(blocksize)
    except OSError:
        return False
    if not chunk:
        return True                       # empty file — treat as (empty) text
    if b"\x00" in chunk:
        return False                      # NUL byte → binary
    try:
        chunk.decode("utf-8")
        return True
    except UnicodeDecodeError:
        # Allow mostly-printable single-byte encodings (latin-1 etc.).
        printable = sum(1 for b in chunk
                        if b in (9, 10, 13) or 32 <= b <= 126 or b >= 128)
        return printable / len(chunk) > 0.85


def describe_binary(path: str) -> str:
    """A compact placeholder for a file whose bytes aren't text or an image."""
    ext = os.path.splitext(path)[1].lower() or "(none)"
    try:
        size = human_size(os.path.getsize(path))
    except OSError:
        size = "unknown size"
    return (f"[binary file — not text/image; cannot be read as content. "
            f"name: {os.path.basename(path)}, type: {ext}, size: {size}]")


def process_file(path: str):
    """Return (label, text_or_None, [image_urls]). Handles ANY file type."""
    ext = os.path.splitext(path)[1].lower()
    name = os.path.basename(path)
    if ext in IMAGE_EXTS:
        return name, None, [image_file_to_url(path)]
    if ext == ".pdf":
        text, images = extract_pdf(path)
        return name, text or None, images
    if ext == ".docx":
        return name, extract_docx(path), []
    if ext == ".pptx":
        return name, extract_pptx(path), []
    if ext in (".xlsx", ".xlsm"):
        try:
            return name, extract_xlsx(path), []
        except Exception as e:
            return name, f"[spreadsheet read failed: {e}]", []
    if ext in TEXT_EXTS or ext == "" or looks_like_text(path):
        return name, read_text_file(path), []
    # Genuinely binary and unrecognised → send a metadata placeholder, not garbage.
    return name, describe_binary(path), []


# ==================================================================== API call
def build_user_content(prompt: str, files: list):
    """Build the content blocks for one user turn (text + images)."""
    content = []
    doc_sections = []
    images = []
    for path in files:
        name, text, imgs = process_file(path)
        images.extend(imgs)
        if text:
            if len(text) > MAX_TEXT_CHARS:
                text = text[:MAX_TEXT_CHARS] + "\n...[truncated]..."
            doc_sections.append(f"\n\n===== FILE: {name} =====\n{text}")

    text_block = prompt.strip()
    if doc_sections:
        text_block += "\n\nAttached document contents:" + "".join(doc_sections)
    if not text_block:
        text_block = "Describe the attached file(s)."

    content.append({"type": "input_text", "text": text_block})
    for url in images:
        content.append({"type": "input_image", "image_url": url})
    return content


# ---- workspace folder: context block, edit protocol, edit parsing ----------
FOLDER_INSTRUCTIONS = (
    "You are working inside the user's workspace folder. The current contents "
    "of its files are provided below under 'WORKSPACE FILES' and are kept in "
    "sync with disk. When the user asks you to change, create, or fix a file, "
    "emit the full new file content in a fenced block whose info string is "
    "`file:<relative/path>`, e.g.\n"
    "```file:src/app.py\n<entire new file content here>\n```\n"
    "Always output the COMPLETE file, not a diff or snippet. "
    "To DELETE a file, emit an empty fenced block whose info string is "
    "`delete:<relative/path>`, e.g.\n"
    "```delete:old.txt\n```\n"
    "NEVER use shell commands like `rm`, `del`, `mv`, or `cp` to change files "
    "— they are NOT executed. Only `file:` and `delete:` blocks take effect. "
    "Use forward slashes in paths, relative to the workspace root. Every block "
    "is applied to disk automatically, logged with a change ID, and can be "
    "reverted by the user, so a previous version is always recoverable. Use "
    "normal prose/code blocks (without the `file:`/`delete:` prefix) for "
    "explanations."
)

# matches ```file:path\n<content>``` and ```delete:path``` (in source order)
_FILE_OP_RE = re.compile(r"```(file|delete):([^\n`]+)\n?(.*?)```", re.S)


def parse_file_ops(text: str):
    """Return [(action, rel_path, content_or_None)] proposed in a message.

    action is "write" or "delete"; content is the new file text for writes.
    """
    ops = []
    for m in _FILE_OP_RE.finditer(text):
        kind = m.group(1)
        rel = m.group(2).strip().strip("`").replace("\\", "/").lstrip("/")
        if not rel:
            continue
        if kind == "delete":
            ops.append(("delete", rel, None))
        else:
            body = m.group(3)
            if body.endswith("\n"):
                body = body[:-1]
            ops.append(("write", rel, body))
    return ops


def build_folder_context(folder_files: list, root_dir: str) -> str:
    """Read folder files fresh from disk and render one labelled text block."""
    if not folder_files:
        return ""
    sections, used = [], 0
    for path in folder_files:
        rel = os.path.relpath(path, root_dir).replace("\\", "/")
        ext = os.path.splitext(path)[1].lower()
        if ext in IMAGE_EXTS:
            continue  # images go through the normal image pipeline
        try:
            _, text, _ = process_file(path)
        except Exception as e:
            text = f"[could not read: {e}]"
        if not text:
            continue
        if used + len(text) > MAX_TEXT_CHARS:
            text = text[: max(0, MAX_TEXT_CHARS - used)] + "\n...[truncated]..."
        used += len(text)
        sections.append(f"\n\n----- {rel} -----\n{text}")
        if used >= MAX_TEXT_CHARS:
            sections.append("\n\n[workspace context truncated — too large]")
            break
    if not sections:
        return ""
    return ("WORKSPACE FILES (live contents, kept in sync with disk):"
            + "".join(sections))


# ---- change history: auto-apply edits, log them, allow revert --------------
HISTORY_DIRNAME = ".fugu_history"


def _history_path(folder: str) -> str:
    return os.path.join(folder, HISTORY_DIRNAME, "history.json")


def load_history(folder: str) -> list:
    """Return the list of change records for a folder (newest last)."""
    try:
        with open(_history_path(folder), "r", encoding="utf-8") as fh:
            data = json.load(fh)
        return data if isinstance(data, list) else []
    except (OSError, json.JSONDecodeError):
        return []


def save_history(folder: str, records: list) -> None:
    d = os.path.join(folder, HISTORY_DIRNAME)
    try:
        os.makedirs(d, exist_ok=True)
        with open(_history_path(folder), "w", encoding="utf-8") as fh:
            json.dump(records, fh, ensure_ascii=False, indent=2)
    except OSError:
        pass


def apply_ops(folder: str, ops: list) -> list:
    """Apply write/delete ops to disk, logging each with an ID. Returns records."""
    records = load_history(folder)
    next_id = max((r.get("id", 0) for r in records), default=0) + 1
    applied = []
    for action, rel, content in ops:
        abspath = os.path.join(folder, rel)
        existed = os.path.isfile(abspath)
        before = None
        if existed:
            try:
                before = read_text_file(abspath)
            except Exception:
                before = None
        rec = {"id": next_id, "time": time.strftime("%Y-%m-%d %H:%M:%S"),
               "file": rel, "before": before, "after": content}
        try:
            if action == "delete":
                if not existed:
                    rec["action"] = "skip"        # nothing to delete
                else:
                    os.remove(abspath)
                    rec["action"] = "delete"
                    rec["after"] = None
            else:                                  # write
                os.makedirs(os.path.dirname(abspath) or ".", exist_ok=True)
                with open(abspath, "w", encoding="utf-8") as fh:
                    fh.write(content)
                rec["action"] = "modify" if existed else "create"
        except OSError as e:
            rec["action"] = "error"; rec["error"] = str(e)
        records.append(rec); applied.append(rec); next_id += 1
    save_history(folder, records)
    return applied


def revert_change(folder: str, change_id: int):
    """Restore the file to its state before change `change_id`. Returns (ok, msg)."""
    records = load_history(folder)
    rec = next((r for r in records if r.get("id") == change_id), None)
    if rec is None:
        return False, f"#{change_id} not found"
    if rec.get("action") not in ("modify", "create", "delete", "revert"):
        return False, f"#{change_id} is not revertable"
    abspath = os.path.join(folder, rec["file"])
    target = rec.get("before")          # state to restore to
    try:
        if target is None:              # the change created the file -> remove it
            if os.path.isfile(abspath):
                os.remove(abspath)
            restored_to = "(deleted)"
        else:
            os.makedirs(os.path.dirname(abspath) or ".", exist_ok=True)
            with open(abspath, "w", encoding="utf-8") as fh:
                fh.write(target)
            restored_to = "previous content"
    except OSError as e:
        return False, f"revert failed: {e}"
    # Log the revert itself as a new change so it too can be undone (redo).
    new_id = max((r.get("id", 0) for r in records), default=0) + 1
    cur_after = rec.get("after")
    records.append({"id": new_id, "time": time.strftime("%Y-%m-%d %H:%M:%S"),
                    "file": rec["file"], "action": "revert", "reverts": change_id,
                    "before": cur_after, "after": target})
    save_history(folder, records)
    return True, f"Reverted #{change_id} ({rec['file']}) → {restored_to}. New entry #{new_id}."


def call_fugu(model: str, history: list, effort: str, instructions: str = None):
    """Send the full conversation history. Returns (text, usage_dict, served_model)."""
    from openai import OpenAI
    client = OpenAI(api_key=API_KEY, base_url=BASE_URL)
    kwargs = {"model": model, "input": history}
    if instructions:
        kwargs["instructions"] = instructions
    if effort and effort != "default":
        kwargs["reasoning"] = {"effort": effort}
    try:
        resp = client.responses.create(**kwargs, timeout=300.0)
    except TypeError:
        kwargs.pop("reasoning", None)
        resp = client.responses.create(**kwargs)
    usage = {}
    try:
        usage = resp.usage.model_dump() if resp.usage else {}
    except Exception:
        usage = {}
    served_model = getattr(resp, "model", None) or model
    return resp.output_text, usage, served_model


# ============================================================== usage tracking
USAGE_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "usage_stats.json")
USAGE_KEYS = ["input_tokens", "output_tokens", "cached_input_tokens",
              "orchestration_input_tokens", "orchestration_output_tokens",
              "orchestration_input_cached_tokens", "reasoning_tokens",
              "total_tokens", "requests"]


def load_usage() -> dict:
    try:
        with open(USAGE_FILE, "r", encoding="utf-8") as fh:
            data = json.load(fh)
    except (OSError, json.JSONDecodeError):
        data = {}
    return {k: int(data.get(k, 0)) for k in USAGE_KEYS}


def add_usage(totals: dict, usage: dict) -> dict:
    """Fold one response's usage dict into the running totals."""
    idet = usage.get("input_tokens_details", {}) or {}
    odet = usage.get("output_tokens_details", {}) or {}
    totals["input_tokens"] += usage.get("input_tokens", 0)
    totals["output_tokens"] += usage.get("output_tokens", 0)
    totals["total_tokens"] += usage.get("total_tokens", 0)
    totals["cached_input_tokens"] += idet.get("cached_tokens", 0)
    totals["orchestration_input_tokens"] += idet.get("orchestration_input_tokens", 0)
    totals["orchestration_input_cached_tokens"] += idet.get("orchestration_input_cached_tokens", 0)
    totals["orchestration_output_tokens"] += odet.get("orchestration_output_tokens", 0)
    totals["reasoning_tokens"] += odet.get("reasoning_tokens", 0)
    totals["requests"] += 1
    try:
        with open(USAGE_FILE, "w", encoding="utf-8") as fh:
            json.dump(totals, fh, indent=2)
    except OSError:
        pass
    return totals


# ============================================================== chat storage
HERE_DIR = os.path.dirname(os.path.abspath(__file__))
CHATS_DIR = os.path.join(HERE_DIR, "chats")


def sanitize_title(title: str) -> str:
    title = re.sub(r'[\\/:*?"<>|\r\n\t]+', " ", title).strip()
    title = re.sub(r"\s+", " ", title)
    return (title or "Untitled")[:60]


def make_title(history: list) -> str:
    """Derive a readable title from the first user message."""
    for msg in history:
        if msg.get("role") == "user":
            for part in msg.get("content", []):
                if part.get("type") == "input_text":
                    raw = part["text"].split("\n\nAttached document contents:")[0]
                    raw = raw.strip().splitlines()[0] if raw.strip() else ""
                    if raw:
                        return sanitize_title(raw)
    return "Untitled chat"


def next_serial() -> int:
    os.makedirs(CHATS_DIR, exist_ok=True)
    nums = []
    for fn in os.listdir(CHATS_DIR):
        m = re.match(r"(\d+)\s*-", fn)
        if m:
            nums.append(int(m.group(1)))
    return (max(nums) + 1) if nums else 1


def save_chat(history: list, display: list, model: str,
              existing_path: str = None, serial: int = None, title: str = None,
              folder: str = None) -> tuple:
    """Write the chat to chats/NNN - Title.json. Returns (path, serial, title)."""
    os.makedirs(CHATS_DIR, exist_ok=True)
    if serial is None:
        serial = next_serial()
    if title is None:
        title = make_title(history)
    data = {
        "serial": serial,
        "title": title,
        "model": model,
        "folder": folder,
        "updated": time.strftime("%Y-%m-%d %H:%M:%S"),
        "history": history,
        "display": display,
    }
    path = os.path.join(CHATS_DIR, f"{serial:03d} - {sanitize_title(title)}.json")
    # If the title changed, remove the old file so we don't leave duplicates.
    if existing_path and os.path.abspath(existing_path) != os.path.abspath(path):
        try:
            os.remove(existing_path)
        except OSError:
            pass
    with open(path, "w", encoding="utf-8") as fh:
        json.dump(data, fh, ensure_ascii=False, indent=2)
    return path, serial, title


def list_chats() -> list:
    """Return [(serial, title, updated, path)] sorted by serial."""
    os.makedirs(CHATS_DIR, exist_ok=True)
    out = []
    for fn in os.listdir(CHATS_DIR):
        if not fn.lower().endswith(".json"):
            continue
        path = os.path.join(CHATS_DIR, fn)
        try:
            with open(path, "r", encoding="utf-8") as fh:
                d = json.load(fh)
            out.append((d.get("serial", 0), d.get("title", fn),
                        d.get("updated", ""), path))
        except (OSError, json.JSONDecodeError):
            continue
    return sorted(out, key=lambda t: t[0])


def load_chat(path: str) -> dict:
    with open(path, "r", encoding="utf-8") as fh:
        return json.load(fh)


# ======================================================================= GUI
def run_gui():
    import tkinter as tk
    from tkinter import ttk, filedialog, scrolledtext

    FILETYPES = [
        ("All files", "*.*"),
        ("Images", "*.png *.jpg *.jpeg *.gif *.webp *.bmp"),
        ("PDF", "*.pdf"), ("Word", "*.docx"), ("PowerPoint", "*.pptx"),
        ("Excel", "*.xlsx *.xlsm"),
        ("Text", "*.txt *.md *.csv *.json"),
    ]
    GREETING = ("Hi! Ask me anything, or attach a file. Each tab is its own "
                "conversation with its own memory — open as many as you like.")

    # Use a drag-and-drop-capable root if tkinterdnd2 is available.
    try:
        from tkinterdnd2 import TkinterDnD, DND_FILES
        root = TkinterDnD.Tk()
        HAS_DND = True
    except Exception:
        root = tk.Tk()
        DND_FILES = None
        HAS_DND = False

    root.title("Fugu Studio — Sakana AI")
    root.geometry("960x820")
    totals = load_usage()   # cumulative token usage, shared across tabs

    # ---- theme palettes (light / dark) ------------------------------------
    ACCENT, ACCENT_HOVER = "#E5402A", "#C8351F"
    PALETTES = {
        "light": dict(bg="#ffffff", panel="#ffffff", fg="#1B1B1B", muted="#777777",
                      text_bg="#ffffff", text_fg="#111111", entry_bg="#ffffff",
                      inline_bg="#eef0f3", head_fg="#0f1b2d", tab_sel="#fbe9e6",
                      sel="#cfe3ff"),
        "dark": dict(bg="#1e1e1e", panel="#252526", fg="#e6edf3", muted="#9aa4b2",
                     text_bg="#1b1f24", text_fg="#e6edf3", entry_bg="#2d2d30",
                     inline_bg="#30363d", head_fg="#cfe0ff", tab_sel="#3a2a28",
                     sel="#264f78"),
    }
    pal = dict(PALETTES["light"])   # live palette read by widgets as they are built
    theme_mode = {"v": "light"}

    style = ttk.Style()
    try:
        style.theme_use("clam")
    except tk.TclError:
        pass

    def apply_theme(mode):
        pal.clear(); pal.update(PALETTES[mode]); theme_mode["v"] = mode
        root.configure(bg=pal["bg"])
        style.configure(".", background=pal["panel"], foreground=pal["fg"])
        style.configure("TFrame", background=pal["panel"])
        style.configure("TLabelframe", background=pal["panel"])
        style.configure("TLabelframe.Label", background=pal["panel"], foreground=pal["muted"])
        style.configure("TLabel", background=pal["panel"], foreground=pal["fg"])
        style.configure("TButton", padding=6)
        style.configure("Accent.TButton", foreground="white", background=ACCENT,
                        padding=8, font=("Segoe UI", 10, "bold"))
        style.map("Accent.TButton", background=[("active", ACCENT_HOVER)])
        style.configure("Title.TLabel", font=("Segoe UI", 18, "bold"),
                        background=pal["panel"], foreground=pal["fg"])
        style.configure("Sub.TLabel", font=("Segoe UI", 10),
                        background=pal["panel"], foreground=pal["muted"])
        style.configure("Credit.TLabel", font=("Segoe UI", 9),
                        background=pal["panel"], foreground=pal["muted"])
        style.configure("TButton", foreground=pal["fg"], background=pal["panel"])
        style.map("TButton",
                  background=[("active", pal["tab_sel"]), ("pressed", pal["tab_sel"])],
                  foreground=[("disabled", pal["muted"])])
        style.configure("TCheckbutton", background=pal["panel"], foreground=pal["fg"])
        style.map("TCheckbutton",
                  background=[("active", pal["panel"])],
                  foreground=[("disabled", pal["muted"])])
        style.configure("TEntry", fieldbackground=pal["entry_bg"], foreground=pal["fg"],
                        insertcolor=pal["fg"])
        style.configure("TCombobox", fieldbackground=pal["entry_bg"],
                        background=pal["entry_bg"], foreground=pal["fg"],
                        arrowcolor=pal["fg"])
        # readonly is the state our comboboxes are in — must be mapped explicitly
        style.map("TCombobox",
                  fieldbackground=[("readonly", pal["entry_bg"]),
                                   ("disabled", pal["panel"])],
                  foreground=[("readonly", pal["fg"]), ("disabled", pal["muted"])],
                  background=[("readonly", pal["entry_bg"])],
                  selectbackground=[("readonly", pal["entry_bg"])],
                  selectforeground=[("readonly", pal["fg"])],
                  arrowcolor=[("!disabled", pal["fg"])])
        # the combobox drop-down list is a classic tk Listbox, themed via options
        root.option_add("*TCombobox*Listbox.background", pal["entry_bg"])
        root.option_add("*TCombobox*Listbox.foreground", pal["fg"])
        root.option_add("*TCombobox*Listbox.selectBackground", pal["sel"])
        root.option_add("*TCombobox*Listbox.selectForeground", pal["fg"])
        # workspace file tree + activity-log table
        style.configure("Treeview", background=pal["entry_bg"],
                        fieldbackground=pal["entry_bg"], foreground=pal["fg"],
                        borderwidth=0)
        style.map("Treeview",
                  background=[("selected", pal["sel"])],
                  foreground=[("selected", pal["fg"])])
        style.configure("Treeview.Heading", background=pal["panel"],
                        foreground=pal["fg"], relief="flat")
        style.map("Treeview.Heading", background=[("active", pal["tab_sel"])])
        style.configure("TScrollbar", background=pal["panel"],
                        troughcolor=pal["bg"], arrowcolor=pal["fg"])
        style.configure("TNotebook", background=pal["bg"], borderwidth=0)
        style.configure("TNotebook.Tab", padding=(12, 6), font=("Segoe UI", 9),
                        background=pal["panel"], foreground=pal["fg"])
        style.map("TNotebook.Tab",
                  background=[("selected", pal["tab_sel"])],
                  foreground=[("selected", ACCENT)])
        # recolor non-ttk widgets that already exist
        try:
            key_dot.configure(bg=pal["panel"])
            usage_lbl.configure(foreground=pal["muted"])
        except Exception:
            pass
        try:
            for ct in list(tabs.values()):
                ct.apply_theme()
        except NameError:
            pass

    root._imgs = {}  # keep PhotoImage refs alive

    def load_logo():
        for name in ("logo.png", "logo.jpg", "logo.jpeg", "fugu_studio.png", "fugu.png"):
            path = os.path.join(_HERE, name)
            if os.path.exists(path):
                try:
                    from PIL import Image, ImageTk
                    img = Image.open(path).convert("RGBA")
                    h = img.copy(); h.thumbnail((220, 56))
                    ic = img.copy(); ic.thumbnail((64, 64))
                    return ImageTk.PhotoImage(ic), ImageTk.PhotoImage(h)
                except Exception:
                    return None, None
        return None, None

    def show_about():
        win = tk.Toplevel(root); win.title("About Fugu Studio")
        win.configure(bg=pal["panel"]); win.geometry("420x300"); win.resizable(False, False)
        if root._imgs.get("header"):
            ttk.Label(win, image=root._imgs["header"]).pack(pady=(18, 6))
        ttk.Label(win, text="Fugu Studio", style="Title.TLabel").pack()
        ttk.Label(win, text="A desktop client for Sakana AI · Fugu & Fugu Ultra",
                  style="Sub.TLabel").pack(pady=(2, 14))
        ttk.Label(win, text="Version 1.1", style="Credit.TLabel").pack()
        ttk.Label(win, text="Created by", style="Credit.TLabel").pack(pady=(10, 0))
        ttk.Label(win, text="Md. Sabbir Ahmed", font=("Segoe UI", 13, "bold"),
                  background="white", foreground=ACCENT).pack()
        ttk.Button(win, text="Close", command=win.destroy).pack(pady=18)
        win.transient(root); win.grab_set()

    # ---- header: logo + credit --------------------------------------------
    header = ttk.Frame(root, padding=(14, 12, 14, 6)); header.pack(fill="x")
    icon_img, header_img = load_logo()
    if header_img:
        root._imgs["icon"], root._imgs["header"] = icon_img, header_img
        try:
            root.iconphoto(True, icon_img)
        except Exception:
            pass
        ttk.Label(header, image=header_img).pack(side="left")
    else:
        ttk.Label(header, text="🐡", font=("Segoe UI", 30)).pack(side="left")
        box = ttk.Frame(header); box.pack(side="left", padx=12)
        ttk.Label(box, text="Fugu Studio", style="Title.TLabel").pack(anchor="w")
        ttk.Label(box, text="sakana ai", style="Sub.TLabel").pack(anchor="w")

    credit = ttk.Frame(header); credit.pack(side="right")
    ttk.Label(credit, text="© MSA", style="Credit.TLabel").pack(anchor="e")
    crow = ttk.Frame(credit); crow.pack(anchor="e", pady=(3, 0))

    def toggle_theme():
        new = "dark" if theme_mode["v"] == "light" else "light"
        apply_theme(new)
        theme_btn.config(text="☀ Light" if new == "dark" else "🌙 Dark")

    theme_btn = ttk.Button(crow, text="🌙 Dark", command=toggle_theme, width=8)
    theme_btn.pack(side="left", padx=(0, 6))
    ttk.Button(crow, text="About", command=show_about, width=8).pack(side="left")

    ttk.Separator(root, orient="horizontal").pack(fill="x", padx=12, pady=(2, 6))

    # ---- API key row (shared) ---------------------------------------------
    keybar = ttk.Frame(root, padding=(10, 0, 10, 4)); keybar.pack(fill="x")
    key_dot = tk.Canvas(keybar, width=14, height=14, highlightthickness=0, bg="white")
    key_dot.pack(side="left", padx=(0, 4))
    _dot = key_dot.create_oval(2, 2, 12, 12, fill="#cf222e", outline="")
    ttk.Label(keybar, text="API key:").pack(side="left")
    key_var = tk.StringVar(value=API_KEY)
    key_state = {"shown": False}
    key_entry = ttk.Entry(keybar, textvariable=key_var, show="*", width=58)
    key_entry.pack(side="left", padx=(4, 4))
    key_status = tk.StringVar(value="")
    key_status_lbl = ttk.Label(keybar, textvariable=key_status)
    key_status_lbl.pack(side="left", padx=6)

    def reflect_key_state():
        if API_KEY:
            key_dot.itemconfig(_dot, fill="#2ea043")
            key_status.set("✓ key set"); key_status_lbl.config(foreground="#2ea043")
        else:
            key_dot.itemconfig(_dot, fill="#cf222e")
            key_status.set("⚠ no API key"); key_status_lbl.config(foreground="#cf222e")

    def toggle_eye():
        key_state["shown"] = not key_state["shown"]
        key_entry.config(show="" if key_state["shown"] else "*")
        eye_btn.config(text="🙈" if key_state["shown"] else "👁")

    def save_key(event=None):
        new = key_var.get().strip()
        if new != API_KEY:
            set_api_key(new)
        reflect_key_state()
        if not key_state["shown"]:
            key_entry.config(show="*")
        return None

    eye_btn = ttk.Button(keybar, text="👁", width=3, command=toggle_eye)
    eye_btn.pack(side="left")
    ttk.Button(keybar, text="Save key", command=save_key).pack(side="left", padx=6)
    key_entry.bind("<FocusIn>", lambda e: key_entry.config(show=""))
    key_entry.bind("<FocusOut>", save_key)
    key_entry.bind("<Return>", lambda e: (root.focus_set(), save_key()))

    # ---- token usage shown on this same row, to the right of the key field
    usage_var = tk.StringVar()
    last_var = tk.StringVar(value="")
    ttk.Button(keybar, text="Reset", width=6,
               command=lambda: reset_usage()).pack(side="right")
    ttk.Label(keybar, textvariable=last_var, foreground="#2ea043",
              font=("Consolas", 9)).pack(side="right", padx=(8, 8))
    usage_lbl = ttk.Label(keybar, textvariable=usage_var, foreground="#444",
                          font=("Consolas", 9))
    usage_lbl.pack(side="right", padx=(20, 0))

    def update_usage_label(last_total=None):
        t = totals
        usage_var.set(
            f"Tokens — in {t['input_tokens']:,}  out {t['output_tokens']:,}  "
            f"cached {t['cached_input_tokens']:,}  |  orch in {t['orchestration_input_tokens']:,}  "
            f"out {t['orchestration_output_tokens']:,}  cached {t['orchestration_input_cached_tokens']:,}  "
            f"|  TOTAL {t['total_tokens']:,}  ({t['requests']} reqs)")
        if last_total:
            last_var.set(f"last +{last_total:,}")

    def reset_usage():
        for k in USAGE_KEYS:
            totals[k] = 0
        try:
            os.remove(USAGE_FILE)
        except OSError:
            pass
        last_var.set(""); update_usage_label()

    # ---- shared helpers ----------------------------------------------------
    _INLINE = re.compile(r"\*\*(.+?)\*\*|__(.+?)__|`([^`]+)`|\*(.+?)\*")
    _FENCE = re.compile(r"```([\w+-]*)\n?(.*?)```", re.S)

    def copy_to_clipboard(text, btn=None, label="Copy"):
        root.clipboard_clear(); root.clipboard_append(text)
        if btn is not None:
            btn.config(text="✓ Copied")
            btn.after(1200, lambda: btn.config(text=label))

    # ---- tab bar + notebook ------------------------------------------------
    tabbar = ttk.Frame(root, padding=(10, 2, 10, 0)); tabbar.pack(fill="x")
    ttk.Button(tabbar, text="＋ New Chat",
               command=lambda: ChatTab()).pack(side="left")
    ttk.Button(tabbar, text="Load chat",
               command=lambda: open_loader()).pack(side="left", padx=6)

    notebook = ttk.Notebook(root)
    notebook.pack(fill="both", expand=True, padx=10, pady=(4, 6))
    tabs = {}  # str(frame) -> ChatTab
    tab_counter = {"n": 0}  # for sequential default names (Chat 1, Chat 2, …)

    # ===================================================== a single chat tab
    class ChatTab:
        def __init__(self, data=None):
            self.history = []
            self.display_log = []
            self.attached = []
            self.folder = None              # selected workspace folder (abs path)
            self.folder_files = []          # abs paths recognised inside it
            self.current = {"path": None, "serial": None, "title": None}
            self.model_var = tk.StringVar(value=(data or {}).get("model", "fugu")
                                          if (data or {}).get("model") in MODELS else "fugu")
            self.effort_var = tk.StringVar(value="default")
            self.status_var = tk.StringVar(value="Ready")
            self.frame = ttk.Frame(notebook)
            self._build()
            if data:
                default_name = f"#{data.get('serial', 0):03d} {data.get('title', 'chat')}"
            else:
                tab_counter["n"] += 1
                default_name = f"Chat {tab_counter['n']}"
            self.name = default_name
            notebook.add(self.frame, text=default_name)
            tabs[str(self.frame)] = self
            notebook.select(self.frame)
            if data:
                self._load_data(data)
            else:
                self.add_bubble("Fugu", GREETING, "fugu", record=False)

        # ---------- UI ----------
        def _build(self):
            f = self.frame
            bar = ttk.Frame(f, padding=(4, 6)); bar.pack(fill="x")
            ttk.Label(bar, text="Model:").pack(side="left")
            ttk.Combobox(bar, textvariable=self.model_var, values=MODELS, width=11,
                         state="readonly").pack(side="left", padx=(4, 14))
            ttk.Label(bar, text="Reasoning:").pack(side="left")
            ttk.Combobox(bar, textvariable=self.effort_var,
                         values=["default", "high", "xhigh"], width=9,
                         state="readonly").pack(side="left", padx=4)
            ttk.Button(bar, text="✕ Close tab", command=self.close).pack(side="right")
            ttk.Button(bar, text="Save", command=self.do_save).pack(side="right", padx=6)
            ttk.Label(bar, textvariable=self.status_var,
                      foreground="#2b6fd6").pack(side="right", padx=8)
            ttk.Button(bar, text="📁 Folder", width=10,
                       command=self.toggle_sidebar).pack(side="left", padx=(10, 0))

            # body = left workspace sidebar + right main (chat) area
            body = ttk.Frame(f); body.pack(side="top", fill="both", expand=True)
            self.sidebar = ttk.Frame(body, width=250, padding=(8, 4))
            self.sidebar_visible = False          # hidden until a folder is browsed
            self.sidebar.pack_propagate(False)
            self._build_sidebar(self.sidebar)
            main = ttk.Frame(body); main.pack(side="left", fill="both", expand=True)

            self.transcript = scrolledtext.ScrolledText(
                main, wrap="word", state="disabled", font=("Segoe UI", 10), padx=8, pady=6,
                height=10,
                background=pal["text_bg"], foreground=pal["text_fg"],
                insertbackground=pal["text_fg"], borderwidth=1, relief="solid")
            t = self.transcript
            t.tag_config("you", foreground="#1f6feb", font=("Segoe UI", 10, "bold"))
            t.tag_config("fugu", foreground="#2ea043", font=("Segoe UI", 10, "bold"))
            t.tag_config("body", foreground=pal["text_fg"])
            t.tag_config("err", foreground="#cf222e")
            t.tag_config("note", foreground="#2ea043", font=("Segoe UI", 9, "bold"))
            t.tag_config("md_b", font=("Segoe UI", 10, "bold"))
            t.tag_config("md_i", font=("Segoe UI", 10, "italic"))
            t.tag_config("md_code", font=("Consolas", 10), background=pal["inline_bg"])
            t.tag_config("md_h", font=("Segoe UI", 13, "bold"), foreground=pal["head_fg"],
                         spacing1=6, spacing3=2)
            t.bind("<Control-c>", self.copy_selection)
            t.bind("<Control-C>", self.copy_selection)

            # Input row and attachment area are packed from the bottom so they
            # always reserve their full height; the transcript fills the rest.
            bottom = ttk.Frame(main, padding=(8, 4, 8, 8)); bottom.pack(side="bottom", fill="x")
            self.prompt_box = tk.Text(bottom, height=4, wrap="word",
                                      font=("Segoe UI", 10),
                                      relief="solid", borderwidth=1,
                                      background=pal["entry_bg"], foreground=pal["fg"],
                                      insertbackground=pal["fg"])
            self.prompt_box.pack(side="left", fill="both", expand=True, ipady=4)
            self.send_btn = ttk.Button(bottom, text="Send  ▶", style="Accent.TButton",
                                       command=self.send, width=10)
            self.send_btn.pack(side="right", padx=(8, 0), ipady=12)

            hint = ("Attach to next message — drag & drop here, paste an image (Ctrl+V), "
                    "or Add files")
            af = ttk.LabelFrame(main, text=hint, padding=6)
            af.pack(side="bottom", fill="x", padx=8, pady=(0, 4))
            self.files_list = tk.Listbox(af, height=2, background=pal["entry_bg"],
                                         foreground=pal["fg"],
                                         selectbackground=pal["sel"])
            self.files_list.pack(side="left", fill="both", expand=True)
            fb = ttk.Frame(af); fb.pack(side="left", padx=8)
            ttk.Button(fb, text="Add files…", command=self.add_files).pack(side="left", padx=1)
            ttk.Button(fb, text="Remove", command=self.remove_selected).pack(side="left", padx=1)
            self.files_list.bind("<Button-3>", lambda e: self.paste())
            if HAS_DND:
                for w in (self.files_list, af):
                    w.drop_target_register(DND_FILES)
                    w.dnd_bind("<<Drop>>", self.on_drop)

            # Transcript packed last → fills the space above the fixed bottom rows.
            self.transcript.pack(fill="both", expand=True, padx=8, pady=(0, 6))
            self.prompt_box.bind("<Return>", self.send)
            self.prompt_box.bind("<Shift-Return>", lambda e: None)

        # ---------- workspace folder sidebar ----------
        def _build_sidebar(self, sb):
            head = ttk.Frame(sb); head.pack(fill="x")
            ttk.Label(head, text="📂 Workspace",
                      font=("Segoe UI", 11, "bold")).pack(side="left")
            ttk.Button(head, text="‹", width=2,
                       command=self.toggle_sidebar).pack(side="right")
            ttk.Button(sb, text="Browse Folder…",
                       command=self.browse_folder).pack(fill="x", pady=(6, 2))
            self.folder_var = tk.StringVar(value="No folder selected")
            ttk.Label(sb, textvariable=self.folder_var, wraplength=224,
                      foreground=pal["muted"], font=("Segoe UI", 8)).pack(fill="x")

            row = ttk.Frame(sb); row.pack(fill="x", pady=(4, 2))
            ttk.Button(row, text="⟳ Rescan", command=self.refresh_folder).pack(side="left")
            ttk.Button(row, text="Clear", command=self.clear_folder).pack(side="left", padx=4)

            self.include_folder = tk.BooleanVar(value=True)
            ttk.Checkbutton(sb, text="Send folder with messages",
                            variable=self.include_folder).pack(fill="x", pady=(2, 2))
            self.folder_stats = tk.StringVar(value="")
            ttk.Label(sb, textvariable=self.folder_stats,
                      foreground=pal["muted"], font=("Segoe UI", 8)).pack(fill="x")

            ttk.Button(sb, text="📜 Activity log",
                       command=self.open_activity_log).pack(fill="x", pady=(6, 1))
            self.log_count_var = tk.StringVar(value="")
            ttk.Label(sb, textvariable=self.log_count_var,
                      foreground=pal["muted"], font=("Segoe UI", 8)).pack(fill="x")

            tree_wrap = ttk.Frame(sb); tree_wrap.pack(fill="both", expand=True, pady=(4, 0))
            self.tree = ttk.Treeview(tree_wrap, show="tree", selectmode="browse")
            ysb = ttk.Scrollbar(tree_wrap, orient="vertical", command=self.tree.yview)
            self.tree.configure(yscrollcommand=ysb.set)
            ysb.pack(side="right", fill="y")
            self.tree.pack(side="left", fill="both", expand=True)
            self.tree.bind("<Double-1>", self.open_tree_item)
            self._tree_paths = {}   # tree item id -> abs path (files only)

        def toggle_sidebar(self):
            if self.sidebar_visible:
                self.sidebar.pack_forget()
            else:
                self.sidebar.pack(side="left", fill="y", before=self.sidebar.master.winfo_children()[-1])
            self.sidebar_visible = not self.sidebar_visible

        def show_sidebar(self):
            if not self.sidebar_visible:
                self.toggle_sidebar()

        def browse_folder(self):
            d = filedialog.askdirectory(title="Select a workspace folder")
            if not d:
                return
            self.folder = os.path.abspath(d)
            self.folder_var.set(self.folder)
            self.show_sidebar()
            self.refresh_folder()

        def refresh_folder(self):
            if not self.folder or not os.path.isdir(self.folder):
                self.folder_stats.set("")
                return
            self.status_var.set("Scanning folder…")
            files, skipped = scan_folder(self.folder)
            self.folder_files = files
            total = sum((os.path.getsize(p) for p in files if os.path.exists(p)), 0)
            note = f"{len(files)} files · {human_size(total)}"
            if skipped:
                note += f" · {skipped} skipped"
            self.folder_stats.set(note)
            self.populate_tree(files)
            self.update_log_count()
            self.status_var.set("Ready")

        def clear_folder(self):
            self.folder = None
            self.folder_files = []
            self.folder_var.set("No folder selected")
            self.folder_stats.set("")
            self.log_count_var.set("")
            self.populate_tree([])

        def populate_tree(self, files):
            self.tree.delete(*self.tree.get_children())
            self._tree_paths = {}
            if not self.folder:
                return
            dir_nodes = {"": ""}   # relative dir -> tree item id ("" == root)
            for path in files:
                rel = os.path.relpath(path, self.folder).replace("\\", "/")
                parts = rel.split("/")
                parent, cur = "", ""
                for d in parts[:-1]:
                    cur = f"{cur}/{d}" if cur else d
                    if cur not in dir_nodes:
                        dir_nodes[cur] = self.tree.insert(
                            dir_nodes[parent], "end", text="📁 " + d, open=False)
                    parent = cur
                node = self.tree.insert(dir_nodes[parent], "end", text="📄 " + parts[-1])
                self._tree_paths[node] = path

        def open_tree_item(self, event):
            item = self.tree.identify_row(event.y)
            path = self._tree_paths.get(item)
            if path and os.path.exists(path):
                try:
                    os.startfile(path)   # Windows: open in default app
                except Exception as e:
                    self.status_var.set(f"Open failed: {e}")

        def apply_theme(self):
            """Recolor this tab's tk widgets to the live palette."""
            self.transcript.configure(background=pal["text_bg"], foreground=pal["text_fg"],
                                      insertbackground=pal["text_fg"])
            self.transcript.tag_config("body", foreground=pal["text_fg"])
            self.transcript.tag_config("md_code", background=pal["inline_bg"])
            self.transcript.tag_config("md_h", foreground=pal["head_fg"])
            self.prompt_box.configure(background=pal["entry_bg"], foreground=pal["fg"],
                                      insertbackground=pal["fg"])
            self.files_list.configure(background=pal["entry_bg"], foreground=pal["fg"],
                                      selectbackground=pal["sel"])

        # ---------- attachments ----------
        def refresh_list(self):
            self.files_list.delete(0, tk.END)
            for p in self.attached:
                self.files_list.insert(tk.END, os.path.basename(p))

        def add_paths(self, paths):
            added = 0
            for p in paths:
                p = p.strip().strip("{}")
                if p and os.path.isfile(p) and p not in self.attached:
                    self.attached.append(p); added += 1
            if added:
                self.refresh_list()
                self.status_var.set(f"Added {added} file(s)")
            return added

        def add_files(self):
            self.add_paths(filedialog.askopenfilenames(title="Choose files",
                                                       filetypes=FILETYPES))

        def remove_selected(self):
            for idx in reversed(self.files_list.curselection()):
                del self.attached[idx]
            self.refresh_list()

        def on_drop(self, event):
            try:
                paths = root.tk.splitlist(event.data)
            except Exception:
                paths = event.data.split()
            self.add_paths(paths)

        def paste(self, event=None):
            try:
                from PIL import ImageGrab
                obj = ImageGrab.grabclipboard()
            except Exception:
                obj = None
            if isinstance(obj, list):
                if self.add_paths(obj):
                    return "break"
            elif obj is not None:
                try:
                    tmp = os.path.join(tempfile.gettempdir(),
                                       f"fugu_paste_{int(time.time()*1000)}.png")
                    obj.convert("RGB").save(tmp, "PNG")
                    self.attached.append(tmp); self.refresh_list()
                    self.status_var.set("Pasted image attached")
                    return "break"
                except Exception as e:
                    self.status_var.set(f"Paste failed: {e}")
            return None

        # ---------- markdown rendering ----------
        def make_code_block(self, code, lang):
            outer = tk.Frame(self.transcript, background="#d7dbe0", bd=0)
            tb = tk.Frame(outer, background="#2b2f36"); tb.pack(fill="x")
            tk.Label(tb, text=(lang or "code"), bg="#2b2f36", fg="#c9d1d9",
                     font=("Segoe UI", 8)).pack(side="left", padx=8, pady=2)
            cbtn = tk.Button(tb, text="Copy", bd=0, bg="#2b2f36", fg="#7ee787",
                             activebackground="#2b2f36", cursor="hand2",
                             font=("Segoe UI", 8, "bold"))
            cbtn.config(command=lambda: copy_to_clipboard(code, cbtn))
            cbtn.pack(side="right", padx=8)
            lines = code.split("\n")
            body = tk.Text(outer, wrap="none", font=("Consolas", 10), bg="#0d1117",
                           fg="#e6edf3", insertbackground="#e6edf3", bd=0,
                           padx=10, pady=8, height=min(len(lines), 28),
                           width=min(max((len(l) for l in lines), default=20), 96))
            body.insert("1.0", code); body.config(state="disabled")
            if len(lines) > 28 or any(len(l) > 96 for l in lines):
                xs = tk.Scrollbar(outer, orient="horizontal", command=body.xview)
                body.config(xscrollcommand=xs.set); body.pack(fill="x"); xs.pack(fill="x")
            else:
                body.pack(fill="x")
            return outer

        def render_inline(self, segment):
            t = self.transcript
            for line in segment.split("\n"):
                stripped = line.lstrip()
                if stripped.startswith("### "):
                    t.insert("end", stripped[4:] + "\n", "md_h"); continue
                if stripped.startswith("## "):
                    t.insert("end", stripped[3:] + "\n", "md_h"); continue
                if stripped.startswith("# "):
                    t.insert("end", stripped[2:] + "\n", "md_h"); continue
                if stripped.startswith(("- ", "* ")):
                    indent = line[:len(line) - len(stripped)]
                    t.insert("end", indent + "•  ", "body")
                    line = stripped[2:]
                pos = 0
                for m in _INLINE.finditer(line):
                    if m.start() > pos:
                        t.insert("end", line[pos:m.start()], "body")
                    bold, under, code, ital = m.groups()
                    if bold is not None:
                        t.insert("end", bold, "md_b")
                    elif under is not None:
                        t.insert("end", under, "md_b")
                    elif code is not None:
                        t.insert("end", code, "md_code")
                    elif ital is not None:
                        t.insert("end", ital, "md_i")
                    pos = m.end()
                t.insert("end", line[pos:] + "\n", "body")

        def render_markdown(self, text):
            t = self.transcript
            pos = 0
            for m in _FENCE.finditer(text):
                self.render_inline(text[pos:m.start()])
                lang, code = m.group(1), m.group(2)
                # File-op fences look like ```file:rel/path or ```delete:rel/path.
                if lang in ("file", "delete") and code.startswith(":"):
                    nl = code.find("\n")
                    rel = (code[1:nl] if nl != -1 else code[1:]).strip()
                    code = code[nl + 1:] if nl != -1 else ""
                    if lang == "delete":
                        self.render_inline(f"🗑 delete {rel}\n")
                        pos = m.end()
                        continue
                    lang = "📄 " + rel
                block = self.make_code_block(code.rstrip("\n"), lang)
                t.window_create("end", window=block, padx=2, pady=4)
                t.insert("end", "\n")
                pos = m.end()
            self.render_inline(text[pos:])

        def add_copy_message_button(self, text):
            btn = tk.Button(self.transcript, text="⧉ Copy", bd=0, fg="#1f6feb",
                            bg=pal["text_bg"], activebackground=pal["text_bg"],
                            cursor="hand2", font=("Segoe UI", 8, "bold"))
            btn.config(command=lambda: copy_to_clipboard(text, btn, "⧉ Copy"))
            self.transcript.window_create("end", window=btn, padx=2)

        def add_bubble(self, who, text, tag, record=True):
            t = self.transcript
            t.config(state="normal")
            if t.index("end-1c") != "1.0":
                t.insert("end", "\n\n")
            t.insert("end", who + "\n", tag)
            if tag == "fugu":
                self.render_markdown(text)
                t.insert("end", "\n")
                self.add_copy_message_button(text)
            else:
                t.insert("end", text, "err" if tag == "err" else "body")
            t.config(state="disabled")
            t.see("end")
            if record:
                self.display_log.append((who, text, tag))

        def copy_selection(self, event=None):
            try:
                sel = self.transcript.get("sel.first", "sel.last")
                if sel:
                    copy_to_clipboard(sel)
            except tk.TclError:
                pass
            return "break"

        # ---------- send / receive ----------
        def set_tab_title(self, text):
            notebook.tab(self.frame, text=(text[:20] + "…") if len(text) > 21 else text)

        def worker(self, model, effort, instructions):
            try:
                result, usage, served = call_fugu(model, list(self.history),
                                                  effort, instructions)
                tag = "fugu"
            except Exception as e:
                result, usage, served = f"⚠️ {e}", {}, model
                tag = "err"
            root.after(0, lambda: self.finish(result, tag, usage, served))

        def finish(self, result, tag, usage, served):
            if tag == "fugu":
                self.history.append({"role": "assistant",
                                     "content": [{"type": "output_text", "text": result}]})
                if usage:
                    add_usage(totals, usage)
                    update_usage_label(usage.get("total_tokens", 0))
            else:
                self.history.pop()
            label = "Error" if tag == "err" else (
                "Fugu Ultra" if served == "fugu-ultra" else "Fugu")
            self.add_bubble(label, result, tag)
            self.status_var.set("Ready")
            self.send_btn.config(state="normal")
            # Auto-apply any file changes Fugu proposed; each is logged with an ID.
            if tag == "fugu" and self.folder:
                ops = parse_file_ops(result)
                if ops:
                    applied = apply_ops(self.folder, ops)
                    self.refresh_folder()
                    ok = [r for r in applied if r["action"] in
                          ("create", "modify", "delete")]
                    if ok:
                        verbs = {"create": "＋", "modify": "✎", "delete": "🗑"}
                        parts = ", ".join(
                            f"{verbs.get(r['action'], '')}{r['file']} [#{r['id']}]"
                            for r in ok)
                        self.add_activity_note(
                            f"✓ Applied {len(ok)} change(s): {parts}")
                        self.status_var.set(f"Applied {len(ok)} change(s)")
                    for r in applied:
                        if r["action"] == "error":
                            self.add_activity_note(
                                f"⚠️ Failed on {r['file']}: {r.get('error','')}",
                                error=True)
                    self.update_log_count()
            if tag == "fugu" and self.current["path"]:
                self.do_save(silent=True)

        def send(self, event=None):
            if not API_KEY:
                self.add_bubble("Error", "⚠️ No API key. Enter your key in the box above "
                                         "and click Save key.", "err", record=False)
                reflect_key_state()
                return "break"
            prompt = self.prompt_box.get("1.0", tk.END).strip()
            if not prompt and not self.attached:
                return "break"
            use_folder = bool(self.folder and self.include_folder.get())
            if use_folder:
                self.refresh_folder()       # live sync: pick up edits/new files
            label = prompt if prompt else "(no text)"
            if self.attached:
                label += "\n📎 " + ", ".join(os.path.basename(p) for p in self.attached)
            if use_folder and self.folder_files:
                label += f"\n📂 {os.path.basename(self.folder)} ({len(self.folder_files)} files)"
            self.add_bubble("You", label, "you")
            content = build_user_content(prompt, list(self.attached))
            if use_folder:
                ctx = build_folder_context(self.folder_files, self.folder)
                if ctx:
                    content.append({"type": "input_text", "text": ctx})
            self.history.append({"role": "user", "content": content})
            self.attached.clear(); self.refresh_list()
            self.prompt_box.delete("1.0", tk.END)
            self.send_btn.config(state="disabled")
            self.status_var.set("Thinking…")
            instructions = FOLDER_INSTRUCTIONS if use_folder else None
            threading.Thread(target=self.worker,
                             args=(self.model_var.get(), self.effort_var.get(),
                                   instructions),
                             daemon=True).start()
            return "break"

        # ---------- activity log / change history ----------
        def add_activity_note(self, msg, error=False):
            t = self.transcript
            t.config(state="normal")
            if t.index("end-1c") != "1.0":
                t.insert("end", "\n")
            t.insert("end", msg + "\n", "err" if error else "note")
            t.config(state="disabled")
            t.see("end")

        def update_log_count(self):
            if not getattr(self, "log_count_var", None):
                return
            if not self.folder:
                self.log_count_var.set("")
                return
            n = len(load_history(self.folder))
            self.log_count_var.set(f"{n} change(s) logged" if n else "no changes yet")

        def open_activity_log(self):
            if not self.folder:
                self.status_var.set("Select a folder first")
                return
            import difflib
            win = tk.Toplevel(root)
            win.title(f"Activity log — {os.path.basename(self.folder)}")
            win.geometry("900x560"); win.configure(bg=pal["panel"])
            ttk.Label(win, text="Every file change Fugu applied. Select one to see "
                                "its diff, then Revert to restore the previous version.",
                      foreground=pal["muted"]).pack(fill="x", padx=10, pady=(8, 4))

            paned = ttk.PanedWindow(win, orient="horizontal")
            paned.pack(fill="both", expand=True, padx=10, pady=4)
            left = ttk.Frame(paned, width=340); left.pack_propagate(False)
            right = ttk.Frame(paned)
            paned.add(left, weight=0); paned.add(right, weight=1)

            cols = ("id", "time", "action", "file")
            tv = ttk.Treeview(left, columns=cols, show="headings", selectmode="browse")
            for c, w in (("id", 44), ("time", 130), ("action", 70), ("file", 150)):
                tv.heading(c, text=c.title()); tv.column(c, width=w, anchor="w")
            ysb = ttk.Scrollbar(left, orient="vertical", command=tv.yview)
            tv.configure(yscrollcommand=ysb.set)
            ysb.pack(side="right", fill="y"); tv.pack(side="left", fill="both", expand=True)

            diff = scrolledtext.ScrolledText(right, wrap="none", state="disabled",
                                             font=("Consolas", 9), background="#0d1117",
                                             foreground="#e6edf3")
            diff.pack(fill="both", expand=True)
            diff.tag_config("add", foreground="#7ee787")
            diff.tag_config("del", foreground="#ff7b72")
            diff.tag_config("hdr", foreground="#9aa4b2")

            recs = {}

            def reload_rows():
                tv.delete(*tv.get_children())
                recs.clear()
                for r in reversed(load_history(self.folder)):   # newest first
                    act = r.get("action")
                    if act == "revert":
                        act = f"revert→#{r.get('reverts','')}"
                    iid = tv.insert("", "end", values=(r.get("id"), r.get("time", ""),
                                                       act, r.get("file", "")))
                    recs[iid] = r
                self.update_log_count()

            def show_sel(event=None):
                sel = tv.selection()
                if not sel:
                    return
                r = recs.get(sel[0])
                if not r:
                    return
                before = (r.get("before") or "").splitlines()
                after = (r.get("after") or "").splitlines()
                lines = difflib.unified_diff(
                    before, after, fromfile=f"{r['file']} (before #{r['id']})",
                    tofile=f"{r['file']} (after #{r['id']})", lineterm="")
                diff.config(state="normal"); diff.delete("1.0", tk.END)
                for ln in lines:
                    tag = ("hdr" if ln.startswith(("+++", "---", "@@")) else
                           "add" if ln.startswith("+") else
                           "del" if ln.startswith("-") else "")
                    diff.insert("end", ln + "\n", tag)
                diff.config(state="disabled")

            tv.bind("<<TreeviewSelect>>", show_sel)

            def do_revert():
                sel = tv.selection()
                if not sel:
                    return
                r = recs.get(sel[0])
                if not r:
                    return
                ok, msg = revert_change(self.folder, r["id"])
                self.status_var.set(msg)
                if ok:
                    self.refresh_folder()
                    reload_rows()

            btns = ttk.Frame(win); btns.pack(fill="x", padx=10, pady=(4, 10))
            ttk.Button(btns, text="↩ Revert selected", style="Accent.TButton",
                       command=do_revert).pack(side="right")
            ttk.Button(btns, text="⟳ Refresh", command=reload_rows).pack(side="right", padx=6)
            ttk.Button(btns, text="Close", command=win.destroy).pack(side="left")
            reload_rows()

        # ---------- persistence ----------
        def do_save(self, silent=False):
            if not self.history:
                if not silent:
                    self.status_var.set("Nothing to save yet")
                return
            path, serial, title = save_chat(
                self.history, self.display_log, self.model_var.get(),
                existing_path=self.current["path"], serial=self.current["serial"],
                title=self.current["title"], folder=self.folder)
            self.current.update(path=path, serial=serial, title=title)
            if not silent:
                self.status_var.set(f"Saved #{serial:03d}")

        def _load_data(self, data):
            self.history = list(data.get("history", []))
            for who, text, tag in data.get("display", []):
                self.add_bubble(who, text, tag)
            self.current.update(path=data.get("_path"), serial=data.get("serial"),
                                title=data.get("title"))
            folder = data.get("folder")
            if folder and os.path.isdir(folder):
                self.folder = os.path.abspath(folder)
                self.folder_var.set(self.folder)
                self.show_sidebar()
                self.refresh_folder()
            self.set_tab_title(f"#{data.get('serial', 0):03d} {data.get('title', 'chat')}")
            self.status_var.set("Loaded")

        def close(self):
            if self.attached:
                self.attached.clear()
            notebook.forget(self.frame)
            tabs.pop(str(self.frame), None)
            if not tabs:               # never leave zero tabs
                ChatTab()

    # ---------- double-click a tab to rename it ----------
    def rename_tab(event):
        try:
            idx = notebook.index(f"@{event.x},{event.y}")
        except tk.TclError:
            return
        tab_id = notebook.tabs()[idx]
        ct = tabs.get(tab_id)
        cur = notebook.tab(tab_id, "text")
        pop = tk.Toplevel(root)
        pop.overrideredirect(True)
        pop.geometry(f"+{notebook.winfo_rootx() + event.x}+{notebook.winfo_rooty() + event.y}")
        var = tk.StringVar(value=cur)
        ent = ttk.Entry(pop, textvariable=var, width=24)
        ent.pack(); ent.focus_force(); ent.select_range(0, "end")

        def commit(e=None):
            name = var.get().strip() or cur
            notebook.tab(tab_id, text=name)
            if ct:
                ct.name = name
                ct.current["title"] = name   # so saving keeps the chosen name
            pop.destroy()

        ent.bind("<Return>", commit)
        ent.bind("<Escape>", lambda e: pop.destroy())
        ent.bind("<FocusOut>", lambda e: pop.destroy())
    notebook.bind("<Double-Button-1>", rename_tab)

    # ---------- global paste routes to the active tab ----------
    def active_tab():
        return tabs.get(notebook.select())

    def paste_clipboard(event=None):
        t = active_tab()
        if t and root.focus_get() is not t.prompt_box:
            return t.paste()
        return None
    root.bind_all("<Control-v>", paste_clipboard)

    # ---------- load a saved chat into a NEW tab ----------
    def open_loader():
        chats = list_chats()
        win = tk.Toplevel(root); win.title("Saved chats"); win.geometry("470x430")
        win.configure(bg=pal["panel"])
        ttk.Label(win, text="Previous chats (open in a new tab):",
                  padding=8).pack(anchor="w")
        lb = tk.Listbox(win, font=("Segoe UI", 10), background=pal["entry_bg"],
                        foreground=pal["fg"], selectbackground=pal["sel"])
        lb.pack(fill="both", expand=True, padx=8)
        for serial, title, updated, _ in chats:
            lb.insert(tk.END, f"#{serial:03d}  {title}    ({updated})")
        if not chats:
            lb.insert(tk.END, "(no saved chats yet)")

        def do_open():
            sel = lb.curselection()
            if not sel or not chats:
                return
            data = load_chat(chats[sel[0]][3])
            data["_path"] = chats[sel[0]][3]
            ChatTab(data)
            win.destroy()

        def do_delete():
            sel = lb.curselection()
            if not sel or not chats:
                return
            try:
                os.remove(chats[sel[0]][3])
            except OSError:
                pass
            win.destroy(); open_loader()

        bar = ttk.Frame(win, padding=8); bar.pack(fill="x")
        ttk.Button(bar, text="Open in new tab", command=do_open).pack(side="left")
        ttk.Button(bar, text="Delete", command=do_delete).pack(side="left", padx=6)
        ttk.Button(bar, text="Cancel", command=win.destroy).pack(side="right")
        lb.bind("<Double-Button-1>", lambda e: do_open())

    # ---- footer credit -----------------------------------------------------
    footer = ttk.Frame(root, padding=(12, 0, 12, 8)); footer.pack(fill="x")
    ttk.Label(footer, text="Fugu Studio · Sakana AI", style="Credit.TLabel").pack(side="left")
    ttk.Label(footer, text="© 2026 Md. Sabbir Ahmed", style="Credit.TLabel").pack(side="right")

    update_usage_label()
    reflect_key_state()
    root._new_tab = lambda data=None: ChatTab(data)   # test/automation hook
    root._tabs = tabs
    root._apply_theme = apply_theme
    ChatTab()              # open the first tab
    apply_theme("light")   # apply ttk styling now that all widgets exist
    root.mainloop()


if __name__ == "__main__":
    try:
        run_gui()
    except Exception:
        # Fallback for environments without a display.
        print(traceback.format_exc())
        input("Press Enter to close...")
